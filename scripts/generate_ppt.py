"""
SMART MATARAM — Generator PPT Karya Inovasi
Jalankan: python scripts/generate_ppt.py
Output  : SMART_MATARAM_Karya_Inovasi.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Warna PLN ──────────────────────────────────────────────────────────────────
TEAL_DARK  = RGBColor(0x00, 0x4D, 0x40)
TEAL       = RGBColor(0x00, 0x89, 0x7B)
TEAL_LIGHT = RGBColor(0xE0, 0xF2, 0xF1)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK       = RGBColor(0x1B, 0x26, 0x31)
GRAY       = RGBColor(0x5D, 0x6D, 0x7E)
YELLOW     = RGBColor(0xFF, 0xD6, 0x00)
RED        = RGBColor(0xC0, 0x39, 0x2B)
GREEN      = RGBColor(0x1E, 0x8B, 0x4C)

# ── Ukuran slide 16:9 ─────────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank = prs.slide_layouts[6]  # fully blank


# ── Helper functions ──────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, size=Pt(14), bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def slide_header(slide, title, subtitle=None, bg=TEAL_DARK):
    """Dark header bar at top"""
    add_rect(slide, 0, 0, W, Inches(1.3), fill_color=bg)
    add_text(slide, title, Inches(0.4), Inches(0.1), Inches(12), Inches(0.7),
             size=Pt(28), bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, Inches(0.4), Inches(0.75), Inches(12), Inches(0.45),
                 size=Pt(13), color=RGBColor(0xB2, 0xDF, 0xDB), align=PP_ALIGN.LEFT)


def bullet_box(slide, title, bullets, x, y, w, h, title_bg=TEAL, card_bg=TEAL_LIGHT,
               title_color=WHITE, text_color=DARK, icon="●"):
    add_rect(slide, x, y, w, Inches(0.42), fill_color=title_bg)
    add_text(slide, title, x + Inches(0.1), y + Inches(0.02), w - Inches(0.2), Inches(0.38),
             size=Pt(12), bold=True, color=title_color)
    add_rect(slide, x, y + Inches(0.42), w, h - Inches(0.42), fill_color=card_bg,
             line_color=TEAL, line_width=Pt(1))
    bullet_text = "\n".join(f"{icon}  {b}" for b in bullets)
    add_text(slide, bullet_text, x + Inches(0.12), y + Inches(0.5),
             w - Inches(0.24), h - Inches(0.6),
             size=Pt(11), color=text_color, wrap=True)


def vs_box(slide, label, items, x, y, w, h, bg):
    add_rect(slide, x, y, w, h, fill_color=bg)
    add_text(slide, label, x, y + Inches(0.05), w, Inches(0.4),
             size=Pt(16), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    text = "\n".join(f"• {i}" for i in items)
    add_text(slide, text, x + Inches(0.15), y + Inches(0.5),
             w - Inches(0.3), h - Inches(0.6),
             size=Pt(11), color=WHITE, wrap=True)


def kpi_box(slide, value, label, x, y, w=Inches(2.8), h=Inches(1.4),
            val_color=TEAL_DARK, bg=TEAL_LIGHT):
    add_rect(slide, x, y, w, h, fill_color=bg, line_color=TEAL, line_width=Pt(1.5))
    add_text(slide, value, x, y + Inches(0.1), w, Inches(0.7),
             size=Pt(28), bold=True, color=val_color, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, y + Inches(0.75), w, Inches(0.55),
             size=Pt(10), color=GRAY, align=PP_ALIGN.CENTER, wrap=True)


def bg_slide(slide, color=RGBColor(0xF4, 0xF6, 0xF8)):
    add_rect(slide, 0, 0, W, H, fill_color=color)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — JUDUL
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank)
add_rect(s1, 0, 0, W, H, fill_color=TEAL_DARK)
add_rect(s1, 0, H - Inches(0.08), W, Inches(0.08), fill_color=YELLOW)

# Ornament circles
for i, (cx, cy, r, alpha) in enumerate([
    (Inches(11.5), Inches(1),   Inches(2.5), RGBColor(0x00,0x69,0x5C)),
    (Inches(12.5), Inches(5.5), Inches(2),   RGBColor(0x00,0x4D,0x40)),
    (Inches(0.5),  Inches(6.5), Inches(1.5), RGBColor(0x00,0x69,0x5C)),
]):
    add_rect(s1, cx - r, cy - r, r*2, r*2, fill_color=alpha)

# Badge "KARYA INOVASI PLN"
add_rect(s1, Inches(0.5), Inches(0.5), Inches(3.5), Inches(0.45), fill_color=YELLOW)
add_text(s1, "KARYA INOVASI PLN", Inches(0.5), Inches(0.5), Inches(3.5), Inches(0.45),
         size=Pt(12), bold=True, color=TEAL_DARK, align=PP_ALIGN.CENTER)

# Nama inovasi
add_text(s1, "SMART", Inches(0.5), Inches(1.4), Inches(12), Inches(1.4),
         size=Pt(80), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s1, "MATARAM", Inches(0.5), Inches(2.6), Inches(12), Inches(1.4),
         size=Pt(80), bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

# Kepanjangan
add_text(s1, "Sistem Monitoring Aset dan Respon Terpadu Mataram",
         Inches(0.5), Inches(3.9), Inches(12), Inches(0.5),
         size=Pt(16), italic=True, color=RGBColor(0xB2, 0xDF, 0xDB), align=PP_ALIGN.CENTER)

# Garis pemisah
add_rect(s1, Inches(3.5), Inches(4.55), Inches(6.3), Inches(0.03), fill_color=YELLOW)

# Tim
add_text(s1, "TIM INOVASI", Inches(0.5), Inches(4.7), Inches(12), Inches(0.35),
         size=Pt(11), bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

tim = [
    "Aruf Budiman  ·  TL OP UP3 Mataram",
    "Lalu Indrawan Saputra  ·  TL Teknik ULP Ampenan",
    "Imam Al Gazali  ·  TL Teknik ULP Tanjung",
]
for i, t in enumerate(tim):
    add_text(s1, t, Inches(0.5), Inches(5.1 + i * 0.35), Inches(12), Inches(0.35),
             size=Pt(12), color=WHITE, align=PP_ALIGN.CENTER)

add_text(s1, "PLN UP3 Mataram  ·  2025",
         Inches(0.5), Inches(6.85), Inches(12), Inches(0.35),
         size=Pt(10), color=GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — MASALAH
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank)
bg_slide(s2)
slide_header(s2, "02  |  MASALAH & KONDISI AWAL",
             "Tantangan yang dihadapi tim operasi PLN UP3 Mataram sebelum inovasi")

masalah = [
    ("🔍  Inspeksi Tidak Termonitor",
     ["Temuan dikirim via WhatsApp Group", "Rawan hilang di antara pesan lain", "Tidak ada tracking status pengerjaan",
      "Tim sulit mencari temuan lama saat mau dikerjakan"]),
    ("⚡  Anomali Gardu Terlambat Dideteksi",
     ["Pengukuran gardu hanya terpantau setelah direkap ke AMG", "Gardu overload & suhu tinggi tidak terdeteksi secara real-time",
      "Respon perbaikan terlambat, risiko kerusakan aset meningkat"]),
    ("📋  Tidak Ada Sistem Evaluasi",
     ["Tidak bisa memantau temuan mana yang sudah/belum dikerjakan", "Tidak ada laporan progres ke atasan secara otomatis",
      "Akuntabilitas tim eksekusi sulit diukur"]),
    ("🔗  Data Terfragmentasi",
     ["Inspeksi jaringan, inspeksi pohon, dan pengukuran gardu di sistem terpisah",
      "Tidak ada satu dashboard tunggal untuk semua informasi", "UP3 tidak bisa memantau kondisi 4 ULP sekaligus"]),
]

cols = [(Inches(0.3), Inches(1.5)), (Inches(6.85), Inches(1.5)),
        (Inches(0.3), Inches(4.3)),  (Inches(6.85), Inches(4.3))]
for (x, y), (title, bullets) in zip(cols, masalah):
    bullet_box(s2, title, bullets, x, y, Inches(6.3), Inches(2.5),
               title_bg=RED, card_bg=RGBColor(0xFF,0xEB,0xEB),
               title_color=WHITE, text_color=DARK, icon="→")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — AKAR MASALAH
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank)
bg_slide(s3)
slide_header(s3, "03  |  ANALISIS AKAR MASALAH",
             "Mengapa kondisi ini bisa terjadi?")

akar = [
    ("Belum Ada Sistem Terintegrasi",
     "Seluruh proses operasi — inspeksi, pengukuran gardu, monitoring gangguan — dilakukan secara manual dan tersebar di platform berbeda (WA, Excel, AMG) tanpa satu sumber data terpusat."),
    ("Proses Manual & Tidak Terstandar",
     "Pelaporan bergantung pada inisiatif individu. Tidak ada workflow baku dari temuan → penugasan → eksekusi → verifikasi selesai, sehingga temuan mudah terlewat."),
    ("Tidak Ada Notifikasi & Alerting",
     "Tidak ada sistem yang otomatis memberi tahu tim saat ada anomali gardu (overload/suhu tinggi) atau temuan inspeksi yang belum ditindak melebihi batas waktu."),
    ("Keterbatasan Visibilitas UP3",
     "Manajer UP3 tidak dapat memantau kondisi seluruh 4 ULP secara bersamaan dalam satu dashboard. Evaluasi hanya bisa dilakukan melalui rapat atau laporan periodik."),
]

for i, (title, desc) in enumerate(akar):
    y = Inches(1.5 + i * 1.38)
    add_rect(s3, Inches(0.3), y, Inches(0.45), Inches(1.1), fill_color=TEAL_DARK)
    add_text(s3, str(i+1), Inches(0.3), y + Inches(0.25), Inches(0.45), Inches(0.6),
             size=Pt(22), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s3, Inches(0.75), y, Inches(12.2), Inches(1.1),
             fill_color=TEAL_LIGHT, line_color=TEAL, line_width=Pt(1))
    add_text(s3, title, Inches(0.9), y + Inches(0.05), Inches(12), Inches(0.38),
             size=Pt(13), bold=True, color=TEAL_DARK)
    add_text(s3, desc, Inches(0.9), y + Inches(0.42), Inches(11.9), Inches(0.6),
             size=Pt(11), color=DARK, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — IDE INOVASI
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank)
bg_slide(s4)
slide_header(s4, "04  |  IDE INOVASI",
             "Satu platform digital untuk semua kebutuhan monitoring operasi distribusi")

# Konsep utama
add_rect(s4, Inches(0.3), Inches(1.5), Inches(12.7), Inches(1.0), fill_color=TEAL_DARK)
add_text(s4, "SMART MATARAM — Web App Monitoring Aset Terpadu Berbasis Cloud",
         Inches(0.5), Inches(1.6), Inches(12.3), Inches(0.8),
         size=Pt(16), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

modul = [
    ("🔍  Inspeksi Jaringan", ["Input temuan langsung dari lapangan", "Tracking status: Temuan → Ditugaskan → Selesai", "Foto sebelum & sesudah perbaikan", "Filter per ULP, penyulang, tim eksekutor"]),
    ("🌳  Inspeksi Pohon/ROW", ["Monitoring pohon berbahaya di jalur penyulang", "Urgency level: Normal, Urgent, Sangat Urgent", "Ditugaskan ke tim PERABASAN", "Peta interaktif lokasi temuan"]),
    ("⚡  Pengukuran Gardu", ["Dashboard beban gardu real-time", "Alert otomatis: overload ≥80%, suhu >60°C", "Work Order PDF otomatis + tracking WO dikirim", "Riwayat pengukuran per gardu"]),
    ("📊  Score Board & Morning Brief", ["Lead Measures mingguan dengan tracking WIN/LOSE", "Rekap gangguan penyulang (SAIDI/SAIFI)", "Morning Brief otomatis ke Telegram jam 08.00 WITA", "Mode presentasi seperti PowerPoint"]),
]

cols2 = [(Inches(0.3), Inches(2.7)), (Inches(3.55), Inches(2.7)),
         (Inches(6.8), Inches(2.7)),  (Inches(10.05), Inches(2.7))]
for (x, y), (title, bullets) in zip(cols2, modul):
    bullet_box(s4, title, bullets, x, y, Inches(3.1), Inches(3.9),
               title_bg=TEAL, card_bg=TEAL_LIGHT, text_color=DARK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — IMPLEMENTASI
# ══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank)
bg_slide(s5)
slide_header(s5, "05  |  IMPLEMENTASI",
             "Bagaimana dan di mana SMART MATARAM diterapkan")

impl = [
    ("🏗️  Teknologi yang Digunakan",
     ["Next.js 16 (web framework modern)", "Supabase (database cloud + autentikasi)", "Google Sheets API (data gangguan existing)", "Vercel (hosting, auto-deploy, cron job)", "Telegram Bot API (notifikasi otomatis)"]),
    ("👥  Sistem Role & Akses",
     ["UP3 — akses semua ULP, semua modul", "Admin ULP — kelola data unit masing-masing", "Inspektor — input temuan lapangan", "HARJAR / HARGAR / PERABASAN — eksekusi temuan", "YANGU / PDKB — task sesuai bidang"]),
    ("📍  Lokasi Penerapan",
     ["ULP Ampenan ✓", "ULP Cakranegara ✓", "ULP Gerung ✓", "ULP Tanjung ✓", "Seluruh ULP di bawah UP3 Mataram sudah aktif"]),
    ("⚙️  Proses Implementasi",
     ["Pengembangan iteratif — fitur aktif langsung dipakai", "Training user per ULP", "Telegram bot aktif kirim laporan tiap pagi", "Integrasi data Google Sheets tanpa migrasi data besar"]),
]

cols3 = [(Inches(0.3), Inches(1.5)), (Inches(3.55), Inches(1.5)),
         (Inches(6.8), Inches(1.5)),  (Inches(10.05), Inches(1.5))]
for (x, y), (title, bullets) in zip(cols3, impl):
    bullet_box(s5, title, bullets, x, y, Inches(3.1), Inches(5.5),
               title_bg=TEAL, card_bg=TEAL_LIGHT, text_color=DARK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — BEFORE vs AFTER
# ══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank)
bg_slide(s6)
slide_header(s6, "06  |  PERBANDINGAN BEFORE vs AFTER",
             "Transformasi nyata dari kondisi lama ke kondisi baru")

before_items = [
    "Temuan inspeksi dikirim via WhatsApp Group",
    "Rawan hilang, tidak ada tracking",
    "Pengukuran gardu hanya terpantau setelah direkap ke AMG",
    "Anomali gardu (overload/suhu) diketahui terlambat",
    "Tidak ada laporan progres otomatis",
    "UP3 harus menunggu laporan manual dari ULP",
    "Work Order tulis tangan / Excel",
    "Data inspeksi, gardu, gangguan di platform terpisah",
]

after_items = [
    "Temuan tercatat di database cloud, tidak hilang",
    "Tracking status real-time: Temuan → Proses → Selesai",
    "Dashboard gardu live, anomali muncul langsung",
    "Alert otomatis overload ≥80% & suhu >60°C",
    "Morning Brief otomatis ke Telegram jam 08.00 WITA",
    "UP3 pantau semua 4 ULP dari satu dashboard",
    "Work Order PDF otomatis, tracking WO terkirim",
    "Satu platform terintegrasi: inspeksi + gardu + gangguan",
]

# VS divider
add_rect(s6, Inches(0.3),  Inches(1.5), Inches(5.9), Inches(5.5), fill_color=RGBColor(0xC0,0x39,0x2B))
add_rect(s6, Inches(7.13), Inches(1.5), Inches(5.9), Inches(5.5), fill_color=TEAL_DARK)
add_rect(s6, Inches(6.0),  Inches(2.5), Inches(1.33), Inches(3.5), fill_color=DARK)

add_text(s6, "SEBELUM", Inches(0.3), Inches(1.55), Inches(5.9), Inches(0.5),
         size=Pt(18), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s6, "SESUDAH", Inches(7.13), Inches(1.55), Inches(5.9), Inches(0.5),
         size=Pt(18), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s6, "VS", Inches(6.0), Inches(3.8), Inches(1.33), Inches(0.9),
         size=Pt(28), bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

before_text = "\n".join(f"✗  {b}" for b in before_items)
after_text  = "\n".join(f"✓  {a}" for a in after_items)

add_text(s6, before_text, Inches(0.5), Inches(2.1), Inches(5.4), Inches(4.5),
         size=Pt(11), color=WHITE, wrap=True)
add_text(s6, after_text, Inches(7.3), Inches(2.1), Inches(5.6), Inches(4.5),
         size=Pt(11), color=WHITE, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — MANFAAT UTAMA
# ══════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank)
bg_slide(s7)
slide_header(s7, "07  |  MANFAAT UTAMA & DAMPAK NYATA",
             "Hasil yang dirasakan langsung oleh tim operasi")

kpis = [
    ("4 ULP", "Unit aktif menggunakan SMART MATARAM"),
    ("100%", "Temuan inspeksi tercatat,\ntidak ada yang hilang"),
    ("Real-time", "Monitoring anomali gardu\noverload & suhu tinggi"),
    ("08.00 WITA", "Morning Brief otomatis\nterkirim ke Telegram tiap hari"),
]
kpi_xs = [Inches(0.3), Inches(3.55), Inches(6.8), Inches(10.05)]
for x, (val, lbl) in zip(kpi_xs, kpis):
    kpi_box(s7, val, lbl, x, Inches(1.5), w=Inches(3.0), h=Inches(1.6))

manfaat = [
    ("📋  Manajemen Inspeksi",
     ["Seluruh temuan dari semua tim (inspektor, mandorline, petugas) terpusat dalam satu sistem",
      "Workflow baku: Temuan → Ditugaskan → Dalam Proses → Selesai",
      "Foto dokumentasi sebelum & sesudah terlampir di setiap temuan",
      "Histori temuan mudah dicari kapan saja untuk evaluasi"]),
    ("⚡  Monitoring Gardu",
     ["Alert otomatis gardu overload (beban ≥80%) dan suhu tinggi (>60°C)",
      "Work Order PDF dibuat otomatis, tidak perlu tulis tangan",
      "Tracking WO sudah terkirim atau belum langsung dari aplikasi",
      "Riwayat pengukuran per gardu tersimpan dan bisa dievaluasi"]),
    ("📊  Pelaporan & Evaluasi",
     ["Score Board Lead Measures dengan status WIN/LOSE tiap minggu",
      "Rekap gangguan penyulang SAIDI & SAIFI otomatis per bulan",
      "Morning Brief harian ke Telegram: rangkuman kondisi semua ULP",
      "Mode Presentasi langsung dari aplikasi untuk rapat manajemen"]),
]

cols4 = [Inches(0.3), Inches(4.55), Inches(8.8)]
for x, (title, bullets) in zip(cols4, manfaat):
    bullet_box(s7, title, bullets, x, Inches(3.3), Inches(4.1), Inches(3.8),
               title_bg=TEAL, card_bg=TEAL_LIGHT, text_color=DARK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — DAMPAK TAMBAHAN
# ══════════════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(blank)
bg_slide(s8)
slide_header(s8, "08  |  DAMPAK TAMBAHAN (MULTIPLIER EFFECT)",
             "Manfaat tidak langsung yang turut dirasakan")

dampak = [
    ("🦺  Keselamatan Kerja (K3)",
     ["Temuan berbahaya (pohon mau roboh, anomali gardu) tidak terlewat",
      "Respon lebih cepat mengurangi risiko kecelakaan kerja",
      "Foto lapangan terdokumentasi sebagai bukti kondisi K3",
      "Pekerjaan berisiko termonitor melalui sistem, bukan asumsi"]),
    ("🌟  Kualitas Pelayanan Pelanggan",
     ["Deteksi anomali gardu lebih cepat → gangguan lebih sedikit",
      "Inspeksi pohon termonitor → risiko SAIDI/SAIFI dari pohon turun",
      "Work Order lebih cepat terbit → respon perbaikan lebih singkat",
      "Potensi penurunan durasi & frekuensi gangguan penyulang"]),
    ("💼  Budaya Kerja & Akuntabilitas",
     ["Tim eksekutor tahu tugasnya dari sistem, bukan hanya WA",
      "Progress pekerjaan tercatat dan bisa dievaluasi oleh atasan",
      "Score Board membangun budaya target mingguan yang terukur",
      "Morning Brief mendorong siklus evaluasi harian yang konsisten"]),
    ("🔄  Efisiensi Operasional",
     ["Tidak perlu rekap manual dari berbagai sumber data",
      "Laporan rapat langsung bisa ditampilkan dari aplikasi",
      "UP3 hemat waktu koordinasi karena data sudah terpusat",
      "Pengambilan keputusan berbasis data, bukan laporan lisan"]),
]

cols5 = [(Inches(0.3), Inches(1.5)), (Inches(3.55), Inches(1.5)),
         (Inches(6.8), Inches(1.5)),  (Inches(10.05), Inches(1.5))]
for (x, y), (title, bullets) in zip(cols5, dampak):
    bullet_box(s8, title, bullets, x, y, Inches(3.1), Inches(5.5),
               title_bg=TEAL_DARK, card_bg=TEAL_LIGHT, text_color=DARK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — REPLIKASI
# ══════════════════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(blank)
bg_slide(s9)
slide_header(s9, "09  |  POTENSI REPLIKASI",
             "Mengapa SMART MATARAM mudah diterapkan di unit lain")

add_rect(s9, Inches(0.3), Inches(1.5), Inches(12.7), Inches(0.8), fill_color=TEAL)
add_text(s9, "Dibangun di atas platform cloud — tidak butuh server fisik, tidak butuh anggaran besar",
         Inches(0.5), Inches(1.55), Inches(12.3), Inches(0.7),
         size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

replikasi = [
    ("☁️  Berbasis Cloud",
     ["Hosting di Vercel (gratis untuk Hobby plan)", "Database Supabase cloud (PostgreSQL terkelola)",
      "Tidak perlu server fisik atau IT infrastructure besar", "Akses dari mana saja via browser"]),
    ("🔧  Mudah Dikonfigurasi",
     ["Sistem role & unit mudah disesuaikan untuk UP3 manapun", "Nama ULP, feeder, penyulang bisa dikonfigurasi",
      "Google Sheets integration pakai spreadsheet existing", "Tidak perlu migrasi data besar"]),
    ("📈  Skalabilitas Tinggi",
     ["Sudah berjalan di 4 ULP sekaligus dalam 1 sistem", "Bisa diperluas ke seluruh UP3 se-NTB",
      "Potensi dikembangkan untuk UP3 nasional", "Arsitektur multi-tenant siap diperluas"]),
    ("🚀  Pengembangan Lanjutan",
     ["IoT sensor gardu untuk data otomatis tanpa input manual", "AI prediksi anomali sebelum terjadi gangguan",
      "Integrasi P2TL dan data pelanggan", "Mobile app native (Android/iOS)"]),
]

cols6 = [(Inches(0.3), Inches(2.5)), (Inches(3.55), Inches(2.5)),
         (Inches(6.8), Inches(2.5)),  (Inches(10.05), Inches(2.5))]
for (x, y), (title, bullets) in zip(cols6, replikasi):
    bullet_box(s9, title, bullets, x, y, Inches(3.1), Inches(4.5),
               title_bg=TEAL, card_bg=TEAL_LIGHT, text_color=DARK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — PENUTUP
# ══════════════════════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(blank)
add_rect(s10, 0, 0, W, H, fill_color=TEAL_DARK)
add_rect(s10, 0, H - Inches(0.08), W, Inches(0.08), fill_color=YELLOW)

add_text(s10, "KESIMPULAN", Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.6),
         size=Pt(13), bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

add_text(s10, "SMART MATARAM", Inches(0.5), Inches(0.9), Inches(12.3), Inches(0.9),
         size=Pt(44), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s10, "Satu Platform. Empat ULP. Semua Termonitor.",
         Inches(0.5), Inches(1.7), Inches(12.3), Inches(0.5),
         size=Pt(16), italic=True, color=RGBColor(0xB2,0xDF,0xDB), align=PP_ALIGN.CENTER)

add_rect(s10, Inches(3.5), Inches(2.3), Inches(6.3), Inches(0.03), fill_color=YELLOW)

kesimpulan = [
    "✅  Menyelesaikan masalah nyata yang dihadapi tim operasi setiap hari",
    "✅  Sudah aktif digunakan oleh 4 ULP di bawah UP3 Mataram",
    "✅  Dibangun dari kebutuhan lapangan, bukan teori",
    "✅  Mudah direplikasi tanpa biaya infrastruktur besar",
    "✅  Terus berkembang: Morning Brief, Score Board, Mode Presentasi",
]
kes_text = "\n".join(kesimpulan)
add_text(s10, kes_text, Inches(1.0), Inches(2.5), Inches(11.3), Inches(2.8),
         size=Pt(13), color=WHITE, wrap=True)

add_rect(s10, Inches(0.3), Inches(5.3), Inches(12.7), Inches(0.8), fill_color=TEAL)
add_text(s10, "🚀  Potensi ke Depan: IoT Sensor Gardu  ·  AI Prediksi Gangguan  ·  Integrasi P2TL  ·  Replikasi Nasional",
         Inches(0.5), Inches(5.35), Inches(12.3), Inches(0.7),
         size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

tim_footer = "Tim Inovasi: Aruf Budiman  ·  Lalu Indrawan Saputra  ·  Imam Al Gazali  |  PLN UP3 Mataram 2025"
add_text(s10, tim_footer, Inches(0.3), Inches(6.3), Inches(12.7), Inches(0.5),
         size=Pt(10), color=GRAY, align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(__file__), "..", "SMART_MATARAM_Karya_Inovasi.pptx")
out_path = os.path.normpath(out_path)
prs.save(out_path)
print(f"OK  File saved: {out_path}")
